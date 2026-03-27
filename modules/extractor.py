"""
ZYN Capital — Modulo de Extracao de Documentos
Utiliza Claude Sonnet API para classificar e extrair dados estruturados
de documentos financeiros (balancos, DREs, matriculas, contratos, etc.).

v2 — Melhorias:
  - Classificacao + extracao unificadas em 1 chamada API
  - Processamento paralelo com ThreadPoolExecutor
  - Cache por SHA256
  - OCR fallback para PDFs escaneados
  - Validacao de CNPJ
"""

import base64
import concurrent.futures
import hashlib
import json
import logging
import mimetypes
import os
import re
import threading
import time
from io import BytesIO
from pathlib import Path
from typing import Any

import anthropic

logger = logging.getLogger(__name__)
# Ensure logs are visible in Streamlit Cloud
if not logger.handlers:
    handler = logging.StreamHandler()
    handler.setFormatter(logging.Formatter("[%(levelname)s] %(name)s: %(message)s"))
    logger.addHandler(handler)
    logger.setLevel(logging.INFO)

MODEL = "claude-haiku-4-5-20251001"  # Haiku para extração — 3x mais rápido, suficiente para classificar e extrair texto
API_DELAY_SECONDS = 0.5  # Delay reduzido — Haiku tem rate limits mais generosos

try:
    CACHE_DIR = Path("/tmp/zyn-credit-engine/output/extraction_cache")
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
except OSError:
    CACHE_DIR = Path(__file__).resolve().parent.parent / "output" / "extraction_cache"
    CACHE_DIR.mkdir(parents=True, exist_ok=True)

# Lock para serializar delays entre chamadas API em threads diferentes
_api_lock = threading.Lock()
_last_api_call = 0.0

TIPOS_DOCUMENTO = [
    "balanco",
    "dre",
    "balancete",
    "demonstracoes_financeiras",
    "matricula",
    "contrato",
    "certidao",
    "ccir_car",
    "laudo_avaliacao",
    "irpf",
    "faturamento",
    "planejamento",
    "cnpj",
    "alteracao_contratual",
    "outro",
]

SCHEMA_BALANCO = {
    "data_base": "",
    "ativo_total": 0,
    "ativo_circulante": 0,
    "caixa_equivalentes": 0,
    "estoques": 0,
    "contas_receber": 0,
    "ativo_nao_circulante": 0,
    "imobilizado": 0,
    "passivo_circulante": 0,
    "emprestimos_cp": 0,
    "fornecedores": 0,
    "passivo_nao_circulante": 0,
    "emprestimos_lp": 0,
    "patrimonio_liquido": 0,
    "capital_social": 0,
}

SCHEMA_DRE = {
    "periodo": "",
    "receita_liquida": 0,
    "custo_mercadorias": 0,
    "lucro_bruto": 0,
    "despesas_operacionais": 0,
    "ebitda": 0,
    "resultado_financeiro": 0,
    "lucro_liquido": 0,
    "margem_bruta_pct": 0,
    "margem_ebitda_pct": 0,
    "margem_liquida_pct": 0,
}

SCHEMA_MATRICULA = {
    "numero_matricula": "",
    "cartorio": "",
    "municipio": "",
    "uf": "",
    "area_ha": 0,
    "proprietario": "",
    "onus": [{"tipo": "", "credor": "", "valor": 0, "data": ""}],
    "averbacoes": [""],
}

SCHEMA_CONTRATO = {
    "tipo_contrato": "",
    "partes": [],
    "objeto": "",
    "valor": 0,
    "prazo": "",
    "garantias": [],
    "clausulas_relevantes": [],
}


# ---------------------------------------------------------------------------
# Helpers: file type detection
# ---------------------------------------------------------------------------

def _get_client() -> anthropic.Anthropic:
    """Retorna cliente Anthropic autenticado via variavel de ambiente."""
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise EnvironmentError(
            "Variavel de ambiente ANTHROPIC_API_KEY nao definida. "
            "Configure-a antes de utilizar o modulo de extracao."
        )
    return anthropic.Anthropic(api_key=api_key, timeout=120)


def _get_media_type(filename: str) -> str:
    """Infere o media type a partir do nome do arquivo."""
    ext = os.path.splitext(filename)[1].lower()
    mapping = {
        ".pdf": "application/pdf",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
    }
    return mapping.get(ext, mimetypes.guess_type(filename)[0] or "application/octet-stream")


def _is_image(filename: str) -> bool:
    """Verifica se o arquivo e uma imagem suportada pela API."""
    ext = os.path.splitext(filename)[1].lower()
    return ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}


def _is_pdf(filename: str) -> bool:
    """Verifica se o arquivo e PDF."""
    return os.path.splitext(filename)[1].lower() == ".pdf"


def _is_xlsx(filename: str) -> bool:
    """Verifica se o arquivo e uma planilha Excel."""
    return os.path.splitext(filename)[1].lower() in {".xlsx", ".xls"}


def _is_docx(filename: str) -> bool:
    """Verifica se o arquivo e um Word .docx."""
    return os.path.splitext(filename)[1].lower() == ".docx"


def _is_pptx(filename: str) -> bool:
    """Verifica se o arquivo e um PowerPoint .pptx."""
    return os.path.splitext(filename)[1].lower() == ".pptx"


# ---------------------------------------------------------------------------
# Helpers: text extraction from various formats
# ---------------------------------------------------------------------------

def _xlsx_to_text(file_bytes: bytes) -> str:
    """Converte arquivo Excel para representacao textual."""
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "Pacote openpyxl necessario para processar arquivos Excel. "
            "Instale com: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Aba: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))

    wb.close()
    return "\n".join(parts)


def _docx_to_text(file_bytes: bytes) -> str:
    """Extrai texto de um arquivo .docx usando python-docx."""
    try:
        from docx import Document
        doc = Document(BytesIO(file_bytes))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Também extrai texto de tabelas
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        logger.warning("Falha ao extrair texto do DOCX: %s", e)
        return ""


def _pptx_to_text(file_bytes: bytes) -> str:
    """Extrai texto de um arquivo .pptx usando python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(file_bytes))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as e:
        logger.warning("Falha ao extrair texto do PPTX: %s", e)
        return ""


def _pdf_to_text(file_bytes: bytes) -> str:
    """Extrai texto de um PDF usando PyPDF2."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(BytesIO(file_bytes))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- Página {i + 1} ---\n{text}")
        return "\n\n".join(pages) if pages else ""
    except Exception as e:
        logger.warning("Falha ao extrair texto do PDF via PyPDF2: %s", e)
        return ""


def _pdf_ocr_fallback(file_bytes: bytes, max_pages: int = 5) -> str:
    """Try OCR on PDF pages. Returns extracted text or empty string.

    Limited to max_pages to avoid hanging on large scanned PDFs.
    Uses 150 DPI (faster) instead of 300 DPI.
    """
    # Strategy 1: pdf2image + pytesseract (best quality)
    try:
        from pdf2image import convert_from_bytes
        import pytesseract

        logger.info("Tentando OCR via pdf2image (max %d páginas, 150 DPI)...", max_pages)
        images = convert_from_bytes(
            file_bytes, dpi=150, first_page=1, last_page=max_pages,
        )
        pages = []
        for i, img in enumerate(images):
            text = pytesseract.image_to_string(img, lang="por")
            if text.strip():
                pages.append(f"--- Página {i + 1} (OCR) ---\n{text}")
        result = "\n\n".join(pages)
        if result.strip():
            logger.info("OCR via pdf2image+pytesseract extraiu %d chars de %d páginas", len(result), len(images))
            return result
    except ImportError:
        logger.debug("pdf2image ou pytesseract nao disponiveis — pulando OCR")
    except Exception as e:
        logger.warning("Falha no OCR via pdf2image+pytesseract: %s", e)

    # Strategy 2: basic Pillow-based approach (render first page from raw bytes)
    try:
        from PIL import Image
        import pytesseract

        img = Image.open(BytesIO(file_bytes))
        text = pytesseract.image_to_string(img, lang="por")
        if text.strip():
            logger.info("OCR via Pillow+pytesseract extraiu %d chars", len(text))
            return text
    except ImportError:
        logger.debug("pytesseract nao disponivel para fallback Pillow")
    except Exception as e:
        logger.debug("Fallback Pillow+pytesseract falhou: %s", e)

    logger.info("OCR nao disponivel ou falhou — PDF será enviado como document block")
    return ""


# ---------------------------------------------------------------------------
# Helpers: content blocks for API
# ---------------------------------------------------------------------------

def _build_content_blocks(file_bytes: bytes, filename: str, text_prompt: str) -> list[dict[str, Any]]:
    """
    Monta os blocos de conteudo para a API do Claude,
    tratando PDFs, imagens e planilhas de forma adequada.

    Para PDFs: extrai texto via PyPDF2 (compatível com todos os planos da API).
    Se o texto for muito curto (scan/imagem), tenta OCR antes de enviar como document block.
    """
    blocks: list[dict[str, Any]] = []

    if _is_xlsx(filename):
        text_content = _xlsx_to_text(file_bytes)
        blocks.append({
            "type": "text",
            "text": (
                f"Conteudo extraido do arquivo '{filename}':\n\n"
                f"{text_content}\n\n---\n\n{text_prompt}"
            ),
        })
    elif _is_docx(filename):
        text_content = _docx_to_text(file_bytes)
        blocks.append({
            "type": "text",
            "text": (
                f"Conteudo extraido do arquivo Word '{filename}':\n\n"
                f"{text_content}\n\n---\n\n{text_prompt}"
            ),
        })
    elif _is_pptx(filename):
        text_content = _pptx_to_text(file_bytes)
        blocks.append({
            "type": "text",
            "text": (
                f"Conteudo extraido da apresentacao '{filename}':\n\n"
                f"{text_content}\n\n---\n\n{text_prompt}"
            ),
        })
    elif _is_image(filename):
        media_type = _get_media_type(filename)
        b64 = base64.standard_b64encode(file_bytes).decode("utf-8")
        blocks.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": media_type,
                "data": b64,
            },
        })
        blocks.append({"type": "text", "text": text_prompt})
    elif _is_pdf(filename):
        # Extrair texto do PDF para enviar como texto puro (mais compatível)
        pdf_text = _pdf_to_text(file_bytes)
        if len(pdf_text.strip()) > 100:
            # Texto suficiente — enviar como texto
            blocks.append({
                "type": "text",
                "text": (
                    f"Conteudo extraido do PDF '{filename}':\n\n"
                    f"{pdf_text}\n\n---\n\n{text_prompt}"
                ),
            })
        else:
            # PDF é scan/imagem — tentar OCR antes de fallback para document block
            ocr_text = _pdf_ocr_fallback(file_bytes)
            if len(ocr_text.strip()) > 100:
                blocks.append({
                    "type": "text",
                    "text": (
                        f"Conteudo extraido via OCR do PDF '{filename}':\n\n"
                        f"{ocr_text}\n\n---\n\n{text_prompt}"
                    ),
                })
            else:
                # Último recurso: enviar como document (requer beta)
                b64 = base64.standard_b64encode(file_bytes).decode("utf-8")
                blocks.append({
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": b64,
                    },
                })
                blocks.append({"type": "text", "text": text_prompt})
    else:
        # Fallback: tenta decodificar como texto
        try:
            text_content = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text_content = file_bytes.decode("latin-1")
        blocks.append({
            "type": "text",
            "text": (
                f"Conteudo do arquivo '{filename}':\n\n"
                f"{text_content}\n\n---\n\n{text_prompt}"
            ),
        })

    return blocks


def _has_document_block(content: list[dict]) -> bool:
    """Verifica se os blocos de conteúdo incluem um document block (PDF scan)."""
    return any(block.get("type") == "document" for block in content)


def _call_api(client: anthropic.Anthropic, content: list[dict], max_tokens: int = 4096, retries: int = 5) -> str:
    """Chama a API Claude com retry, rate-limit handling e throttling global."""
    global _last_api_call

    kwargs = {
        "model": MODEL,
        "max_tokens": max_tokens,
        "messages": [{"role": "user", "content": content}],
    }
    if _has_document_block(content):
        kwargs["extra_headers"] = {"anthropic-beta": "pdfs-2024-09-25"}

    # Throttle: garante intervalo minimo entre chamadas (thread-safe)
    with _api_lock:
        now = time.monotonic()
        elapsed = now - _last_api_call
        if elapsed < API_DELAY_SECONDS:
            time.sleep(API_DELAY_SECONDS - elapsed)
        _last_api_call = time.monotonic()

    for attempt in range(retries):
        try:
            response = client.messages.create(**kwargs)
            return response.content[0].text
        except anthropic.RateLimitError:
            wait = (attempt + 1) * 15  # 15s, 30s, 45s
            logger.warning("Rate limit atingido. Aguardando %ds antes de retry %d/%d", wait, attempt + 1, retries)
            time.sleep(wait)
        except anthropic.BadRequestError as e:
            if "PDF" in str(e) and _has_document_block(content):
                # PDF inválido como document — tenta extrair texto e reenviar
                logger.warning("PDF inválido como document block. Tentando como texto...")
                new_content = []
                for block in content:
                    if block.get("type") == "document":
                        new_content.append({
                            "type": "text",
                            "text": "[PDF não processável — documento escaneado sem texto extraível]",
                        })
                    else:
                        new_content.append(block)
                kwargs["messages"] = [{"role": "user", "content": new_content}]
                if "extra_headers" in kwargs:
                    del kwargs["extra_headers"]
                response = client.messages.create(**kwargs)
                return response.content[0].text
            raise

    raise Exception("Rate limit excedido após todas as tentativas. Tente novamente em alguns minutos.")


def _parse_json_response(text: str) -> dict:
    """
    Tenta extrair JSON da resposta do Claude.
    Primeiro tenta json.loads direto; se falhar, busca bloco JSON via regex.
    """
    # Tenta parse direto (resposta e JSON puro)
    cleaned = text.strip()
    try:
        return json.loads(cleaned)
    except (json.JSONDecodeError, ValueError):
        pass

    # Tenta encontrar bloco JSON em markdown (```json ... ```)
    match = re.search(r"```(?:json)?\s*\n?(.*?)\n?\s*```", cleaned, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1).strip())
        except (json.JSONDecodeError, ValueError):
            pass

    # Tenta encontrar qualquer objeto JSON na resposta
    match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(0))
        except (json.JSONDecodeError, ValueError):
            pass

    logger.warning("Nao foi possivel extrair JSON da resposta do modelo.")
    return {"error": "Falha ao interpretar resposta do modelo", "resposta_bruta": text}


# ---------------------------------------------------------------------------
# Cache por SHA256
# ---------------------------------------------------------------------------

def _get_file_hash(file_bytes: bytes) -> str:
    """Retorna hash SHA256 truncado (16 hex chars) do conteudo do arquivo."""
    return hashlib.sha256(file_bytes).hexdigest()[:16]


def _get_cached(file_hash: str) -> dict | None:
    """Busca resultado em cache pelo hash. Retorna None se nao encontrado."""
    cache_file = CACHE_DIR / f"{file_hash}.json"
    if cache_file.exists():
        try:
            with open(cache_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            logger.info("Cache hit para hash %s", file_hash)
            return data
        except (json.JSONDecodeError, OSError) as e:
            logger.warning("Cache corrompido para hash %s: %s", file_hash, e)
    return None


def _save_cache(file_hash: str, result: dict):
    """Salva resultado no cache."""
    cache_file = CACHE_DIR / f"{file_hash}.json"
    try:
        with open(cache_file, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        logger.debug("Cache salvo para hash %s", file_hash)
    except OSError as e:
        logger.warning("Falha ao salvar cache para hash %s: %s", file_hash, e)


# ---------------------------------------------------------------------------
# Validacao CNPJ
# ---------------------------------------------------------------------------

def validate_cnpj(cnpj: str) -> bool:
    """Validates CNPJ check digits."""
    # Remove caracteres nao numericos
    cnpj = re.sub(r"\D", "", cnpj)

    if len(cnpj) != 14:
        return False

    # Rejeita CNPJs com todos os digitos iguais
    if cnpj == cnpj[0] * 14:
        return False

    # Calculo do primeiro digito verificador
    pesos_1 = [5, 4, 3, 2, 9, 8, 7, 6, 5, 4, 3, 2]
    soma = sum(int(cnpj[i]) * pesos_1[i] for i in range(12))
    resto = soma % 11
    digito_1 = 0 if resto < 2 else 11 - resto

    if int(cnpj[12]) != digito_1:
        return False

    # Calculo do segundo digito verificador
    pesos_2 = [6, 5, 4, 3, 2, 9, 8, 7, 6, 5, 4, 3, 2]
    soma = sum(int(cnpj[i]) * pesos_2[i] for i in range(13))
    resto = soma % 11
    digito_2 = 0 if resto < 2 else 11 - resto

    if int(cnpj[13]) != digito_2:
        return False

    return True


# ---------------------------------------------------------------------------
# Prompts de extracao por tipo de documento
# ---------------------------------------------------------------------------

def _get_extraction_prompt(tipo_documento: str) -> str:
    """Retorna o prompt de extracao adequado para o tipo de documento."""

    prompts = {
        "balanco": (
            "Voce e um analista financeiro especializado. Extraia todos os dados do Balanco Patrimonial "
            "contidos neste documento. Retorne EXCLUSIVAMENTE um JSON com a seguinte estrutura:\n\n"
            f"{json.dumps(SCHEMA_BALANCO, indent=2, ensure_ascii=False)}\n\n"
            "Regras:\n"
            "- Valores monetarios em reais (R$), sem separador de milhar, ponto como decimal.\n"
            "- Se um campo nao estiver presente, mantenha o valor padrao (0 ou string vazia).\n"
            "- data_base no formato YYYY-MM-DD.\n"
            "- Nao inclua texto adicional fora do JSON."
        ),
        "dre": (
            "Voce e um analista financeiro especializado. Extraia todos os dados da DRE "
            "(Demonstracao de Resultado do Exercicio) contidos neste documento. "
            "Retorne EXCLUSIVAMENTE um JSON com a seguinte estrutura:\n\n"
            f"{json.dumps(SCHEMA_DRE, indent=2, ensure_ascii=False)}\n\n"
            "Regras:\n"
            "- Valores monetarios em reais (R$), sem separador de milhar, ponto como decimal.\n"
            "- Margens em percentual (ex: 25.5 para 25,5%).\n"
            "- periodo no formato 'YYYY' ou 'YYYY-MM a YYYY-MM'.\n"
            "- Se um campo nao estiver presente, mantenha o valor padrao.\n"
            "- Nao inclua texto adicional fora do JSON."
        ),
        "balancete": (
            "Voce e um analista financeiro especializado. Extraia os dados do balancete contabil. "
            "Retorne EXCLUSIVAMENTE um JSON com a seguinte estrutura:\n\n"
            f"{json.dumps(SCHEMA_BALANCO, indent=2, ensure_ascii=False)}\n\n"
            "Regras:\n"
            "- Utilize a mesma estrutura do Balanco Patrimonial, preenchendo os campos disponiveis.\n"
            "- Valores monetarios em reais (R$), sem separador de milhar, ponto como decimal.\n"
            "- Se um campo nao estiver presente no balancete, mantenha o valor padrao.\n"
            "- Nao inclua texto adicional fora do JSON."
        ),
        "matricula": (
            "Voce e um analista juridico especializado em direito imobiliario rural. "
            "Extraia os dados da matricula de imovel contida neste documento. "
            "Retorne EXCLUSIVAMENTE um JSON com a seguinte estrutura:\n\n"
            f"{json.dumps(SCHEMA_MATRICULA, indent=2, ensure_ascii=False)}\n\n"
            "Regras:\n"
            "- area_ha em hectares (numero decimal).\n"
            "- onus: liste todos os onus/gravames encontrados (hipotecas, alienacoes fiduciarias, penhoras, etc.).\n"
            "- averbacoes: liste as averbacoes mais relevantes (reserva legal, construcao, etc.).\n"
            "- Valores monetarios em reais (R$), sem separador de milhar, ponto como decimal.\n"
            "- Datas no formato YYYY-MM-DD.\n"
            "- Nao inclua texto adicional fora do JSON."
        ),
        "contrato": (
            "Voce e um analista juridico especializado em contratos financeiros. "
            "Extraia os dados do contrato contido neste documento. "
            "Retorne EXCLUSIVAMENTE um JSON com a seguinte estrutura:\n\n"
            f"{json.dumps(SCHEMA_CONTRATO, indent=2, ensure_ascii=False)}\n\n"
            "Regras:\n"
            "- partes: liste todas as partes (nomes completos ou razoes sociais).\n"
            "- garantias: liste todas as garantias mencionadas.\n"
            "- clausulas_relevantes: resuma as clausulas mais importantes (vencimento antecipado, cross-default, etc.).\n"
            "- Valor em reais (R$), sem separador de milhar, ponto como decimal.\n"
            "- Nao inclua texto adicional fora do JSON."
        ),
        "certidao": (
            "Voce e um analista de credito. Extraia os dados da certidao contida neste documento. "
            "Retorne EXCLUSIVAMENTE um JSON com os seguintes campos (adapte conforme o tipo de certidao):\n\n"
            '{"tipo_certidao": "", "orgao_emissor": "", "data_emissao": "", "validade": "", '
            '"nome_consultado": "", "cpf_cnpj": "", "resultado": "positiva|negativa|positiva_com_efeito_negativa", '
            '"detalhes": [], "observacoes": ""}\n\n'
            "Regras:\n"
            "- Datas no formato YYYY-MM-DD.\n"
            "- detalhes: liste debitos, protestos ou distribuicoes encontrados, se houver.\n"
            "- Nao inclua texto adicional fora do JSON."
        ),
        "ccir_car": (
            "Voce e um analista especializado em documentacao rural. "
            "Extraia os dados do CCIR ou CAR contido neste documento. "
            "Retorne EXCLUSIVAMENTE um JSON com os seguintes campos:\n\n"
            '{"tipo": "CCIR|CAR", "codigo": "", "nome_imovel": "", "municipio": "", "uf": "", '
            '"area_total_ha": 0, "area_reserva_legal_ha": 0, "area_app_ha": 0, '
            '"proprietario": "", "cpf_cnpj": "", "situacao": "", "data_emissao": ""}\n\n'
            "Regras:\n"
            "- Areas em hectares (numero decimal).\n"
            "- Datas no formato YYYY-MM-DD.\n"
            "- Nao inclua texto adicional fora do JSON."
        ),
    }

    prompts["demonstracoes_financeiras"] = (
        "Voce e um analista financeiro especializado. Este documento contem Demonstracoes Financeiras completas. "
        "Extraia os dados do Balanco Patrimonial E da DRE em um unico JSON.\n\n"
        "Retorne EXCLUSIVAMENTE um JSON com a seguinte estrutura:\n\n"
        '{"data_base": "", "periodo_dre": "",\n'
        ' "ativo_total": 0, "ativo_circulante": 0, "caixa_equivalentes": 0, "estoques": 0, "contas_receber": 0,\n'
        ' "ativo_nao_circulante": 0, "imobilizado": 0,\n'
        ' "passivo_circulante": 0, "emprestimos_cp": 0, "fornecedores": 0,\n'
        ' "passivo_nao_circulante": 0, "emprestimos_lp": 0,\n'
        ' "patrimonio_liquido": 0, "capital_social": 0,\n'
        ' "receita_liquida": 0, "custo_mercadorias": 0, "lucro_bruto": 0,\n'
        ' "despesas_operacionais": 0, "ebitda": 0, "resultado_financeiro": 0, "lucro_liquido": 0,\n'
        ' "margem_bruta_pct": 0, "margem_ebitda_pct": 0, "margem_liquida_pct": 0}\n\n'
        "Regras:\n"
        "- Valores monetarios em numeros puros (sem R$, sem separador de milhar, ponto como decimal).\n"
        "- Se houver dados de mais de um exercicio, use o mais recente.\n"
        "- Margens em percentual (ex: 25.5 para 25,5%).\n"
        "- Responda APENAS com o JSON."
    )
    prompts["laudo_avaliacao"] = (
        "Voce e um analista de credito. Extraia os dados deste Laudo/Relatorio de Avaliacao.\n"
        "Retorne EXCLUSIVAMENTE um JSON:\n\n"
        '{"tipo_laudo": "", "elaborado_por": "", "data_laudo": "",\n'
        ' "imovel_nome": "", "municipio": "", "uf": "", "area_ha": 0,\n'
        ' "valor_mercado": 0, "valor_liquidacao": 0,\n'
        ' "metodo_avaliacao": "", "finalidade": "",\n'
        ' "observacoes": ""}\n\n'
        "Valores monetarios em numeros puros. Responda APENAS com o JSON."
    )
    prompts["irpf"] = (
        "Voce e um analista de credito. Extraia os dados desta Declaracao de IRPF.\n"
        "Retorne EXCLUSIVAMENTE um JSON:\n\n"
        '{"contribuinte": "", "cpf": "", "exercicio": "", "ano_calendario": "",\n'
        ' "rendimentos_tributaveis": 0, "rendimentos_isentos": 0,\n'
        ' "bens_direitos_total": 0, "bens_direitos": [{"descricao": "", "valor": 0}],\n'
        ' "dividas_onus_total": 0, "dividas_onus": [{"descricao": "", "valor": 0}],\n'
        ' "imposto_devido": 0, "atividade_rural_resultado": 0}\n\n'
        "Valores monetarios em numeros puros. Responda APENAS com o JSON."
    )
    prompts["faturamento"] = (
        "Voce e um analista financeiro. Extraia os dados deste relatorio de faturamento/receita.\n"
        "Retorne EXCLUSIVAMENTE um JSON:\n\n"
        '{"empresa": "", "periodo": "", "faturamento_total": 0,\n'
        ' "detalhamento_mensal": [{"mes": "", "valor": 0}],\n'
        ' "detalhamento_empresa": [{"empresa": "", "valor": 0}],\n'
        ' "observacoes": ""}\n\n'
        "Valores monetarios em numeros puros. Responda APENAS com o JSON."
    )
    prompts["planejamento"] = (
        "Voce e um analista financeiro. Extraia os dados deste planejamento/projecao.\n"
        "Retorne EXCLUSIVAMENTE um JSON:\n\n"
        '{"descricao": "", "periodo": "", "dados_projetados": {},\n'
        ' "premissas": [], "observacoes": ""}\n\n'
        "Adapte os campos conforme o conteudo. Responda APENAS com o JSON."
    )
    prompts["cnpj"] = (
        "Extraia os dados deste Cartao CNPJ ou consulta cadastral.\n"
        "Retorne EXCLUSIVAMENTE um JSON:\n\n"
        '{"razao_social": "", "cnpj": "", "nome_fantasia": "", "data_abertura": "",\n'
        ' "natureza_juridica": "", "atividade_principal": "", "endereco": "",\n'
        ' "municipio": "", "uf": "", "situacao_cadastral": "", "data_situacao": "",\n'
        ' "socios": [{"nome": "", "qualificacao": ""}]}\n\n'
        "Responda APENAS com o JSON."
    )
    prompts["alteracao_contratual"] = (
        "Voce e um analista juridico. Extraia os dados desta alteracao contratual / aditivo.\n"
        "Retorne EXCLUSIVAMENTE um JSON:\n\n"
        '{"empresa": "", "cnpj": "", "tipo_alteracao": "", "data": "",\n'
        ' "objeto_alteracao": "", "capital_social": 0,\n'
        ' "socios": [{"nome": "", "participacao_pct": 0}],\n'
        ' "clausulas_relevantes": []}\n\n'
        "Valores monetarios em numeros puros. Responda APENAS com o JSON."
    )

    return prompts.get(
        tipo_documento,
        (
            "Voce e um analista de credito estruturado da ZYN Capital. "
            "Analise o documento e extraia TODAS as informacoes relevantes para analise de credito.\n\n"
            "Retorne EXCLUSIVAMENTE um JSON estruturado com os campos que fizerem sentido para este documento. "
            "Use nomes de campos descritivos em portugues (snake_case). "
            "Valores monetarios como numeros puros. Datas como YYYY-MM-DD.\n\n"
            "Responda APENAS com o JSON, sem texto adicional."
        ),
    )


# ---------------------------------------------------------------------------
# Unified classify + extract prompt
# ---------------------------------------------------------------------------

def _build_unified_prompt(tipo_list: list[str]) -> str:
    """Monta prompt unificado que classifica E extrai em uma unica chamada."""
    categorias = (
        "Categorias possiveis:\n"
        "- balanco: Balanco Patrimonial (isolado)\n"
        "- dre: DRE isolada\n"
        "- balancete: Balancete contabil\n"
        "- demonstracoes_financeiras: Demonstracoes Financeiras completas (BP + DRE + DMPL + DFC + Notas)\n"
        "- matricula: Matricula de imovel rural ou urbano\n"
        "- contrato: Contrato (emprestimo, arrendamento, compra e venda, etc.)\n"
        "- certidao: Certidao (negativa de debitos, protestos, distribuicao, etc.)\n"
        "- ccir_car: CCIR ou CAR (Cadastro Ambiental Rural)\n"
        "- laudo_avaliacao: Laudo ou Relatorio Tecnico de Avaliacao de imoveis/ativos\n"
        "- irpf: Declaracao de Imposto de Renda Pessoa Fisica (DIRPF)\n"
        "- faturamento: Relatorio de faturamento, receita ou analise de receita\n"
        "- planejamento: Planejamento de producao, orcamento ou projecao\n"
        "- cnpj: Cartao CNPJ, comprovante de inscricao ou consulta CNPJ\n"
        "- alteracao_contratual: Alteracao contratual, aditivo ou consolidacao societaria\n"
        "- outro: Documento que nao se encaixe em nenhuma categoria acima\n"
    )

    # Build extraction schemas section
    schemas_section = "\n\nSchemas de extracao por tipo (use o schema correspondente ao tipo classificado):\n\n"
    schemas_section += f'balanco / balancete:\n{json.dumps(SCHEMA_BALANCO, indent=2, ensure_ascii=False)}\n\n'
    schemas_section += f'dre:\n{json.dumps(SCHEMA_DRE, indent=2, ensure_ascii=False)}\n\n'
    schemas_section += f'matricula:\n{json.dumps(SCHEMA_MATRICULA, indent=2, ensure_ascii=False)}\n\n'
    schemas_section += f'contrato:\n{json.dumps(SCHEMA_CONTRATO, indent=2, ensure_ascii=False)}\n\n'
    schemas_section += (
        'demonstracoes_financeiras:\n'
        '{"data_base": "", "periodo_dre": "", "ativo_total": 0, "ativo_circulante": 0, '
        '"caixa_equivalentes": 0, "estoques": 0, "contas_receber": 0, "ativo_nao_circulante": 0, '
        '"imobilizado": 0, "passivo_circulante": 0, "emprestimos_cp": 0, "fornecedores": 0, '
        '"passivo_nao_circulante": 0, "emprestimos_lp": 0, "patrimonio_liquido": 0, "capital_social": 0, '
        '"receita_liquida": 0, "custo_mercadorias": 0, "lucro_bruto": 0, "despesas_operacionais": 0, '
        '"ebitda": 0, "resultado_financeiro": 0, "lucro_liquido": 0, "margem_bruta_pct": 0, '
        '"margem_ebitda_pct": 0, "margem_liquida_pct": 0}\n\n'
    )
    schemas_section += (
        'certidao:\n'
        '{"tipo_certidao": "", "orgao_emissor": "", "data_emissao": "", "validade": "", '
        '"nome_consultado": "", "cpf_cnpj": "", "resultado": "positiva|negativa|positiva_com_efeito_negativa", '
        '"detalhes": [], "observacoes": ""}\n\n'
    )
    schemas_section += (
        'ccir_car:\n'
        '{"tipo": "CCIR|CAR", "codigo": "", "nome_imovel": "", "municipio": "", "uf": "", '
        '"area_total_ha": 0, "area_reserva_legal_ha": 0, "area_app_ha": 0, '
        '"proprietario": "", "cpf_cnpj": "", "situacao": "", "data_emissao": ""}\n\n'
    )
    schemas_section += (
        'laudo_avaliacao:\n'
        '{"tipo_laudo": "", "elaborado_por": "", "data_laudo": "", "imovel_nome": "", '
        '"municipio": "", "uf": "", "area_ha": 0, "valor_mercado": 0, "valor_liquidacao": 0, '
        '"metodo_avaliacao": "", "finalidade": "", "observacoes": ""}\n\n'
    )
    schemas_section += (
        'irpf:\n'
        '{"contribuinte": "", "cpf": "", "exercicio": "", "ano_calendario": "", '
        '"rendimentos_tributaveis": 0, "rendimentos_isentos": 0, "bens_direitos_total": 0, '
        '"bens_direitos": [{"descricao": "", "valor": 0}], "dividas_onus_total": 0, '
        '"dividas_onus": [{"descricao": "", "valor": 0}], "imposto_devido": 0, '
        '"atividade_rural_resultado": 0}\n\n'
    )
    schemas_section += (
        'faturamento:\n'
        '{"empresa": "", "periodo": "", "faturamento_total": 0, '
        '"detalhamento_mensal": [{"mes": "", "valor": 0}], '
        '"detalhamento_empresa": [{"empresa": "", "valor": 0}], "observacoes": ""}\n\n'
    )
    schemas_section += (
        'planejamento:\n'
        '{"descricao": "", "periodo": "", "dados_projetados": {}, "premissas": [], "observacoes": ""}\n\n'
    )
    schemas_section += (
        'cnpj:\n'
        '{"razao_social": "", "cnpj": "", "nome_fantasia": "", "data_abertura": "", '
        '"natureza_juridica": "", "atividade_principal": "", "endereco": "", "municipio": "", '
        '"uf": "", "situacao_cadastral": "", "data_situacao": "", '
        '"socios": [{"nome": "", "qualificacao": ""}]}\n\n'
    )
    schemas_section += (
        'alteracao_contratual:\n'
        '{"empresa": "", "cnpj": "", "tipo_alteracao": "", "data": "", "objeto_alteracao": "", '
        '"capital_social": 0, "socios": [{"nome": "", "participacao_pct": 0}], "clausulas_relevantes": []}\n\n'
    )
    schemas_section += (
        'outro:\n'
        'JSON livre com campos descritivos em portugues (snake_case).\n'
    )

    prompt = (
        "Voce e um analista de credito estruturado da ZYN Capital. "
        "Execute DUAS tarefas em uma unica resposta:\n\n"
        "TAREFA 1 — CLASSIFICACAO: Identifique o tipo do documento.\n"
        f"{categorias}\n"
        "Se o documento contiver BP + DRE juntos, classifique como 'demonstracoes_financeiras'.\n\n"
        "TAREFA 2 — EXTRACAO: Extraia os dados estruturados conforme o schema do tipo identificado.\n"
        f"{schemas_section}\n"
        "Regras gerais de extracao:\n"
        "- Valores monetarios como numeros puros (sem R$, sem separador de milhar, ponto como decimal).\n"
        "- Datas no formato YYYY-MM-DD.\n"
        "- Se um campo nao estiver presente, mantenha o valor padrao (0 ou string vazia).\n"
        "- Margens em percentual (ex: 25.5 para 25,5%).\n\n"
        "Responda EXCLUSIVAMENTE com um JSON no seguinte formato (sem texto adicional):\n"
        '{"classificacao": {"tipo": "<categoria>", "confianca": <0.0 a 1.0>, "descricao": "<breve descricao>"}, '
        '"dados": {<dados extraidos conforme schema do tipo>}}'
    )
    return prompt


# ---------------------------------------------------------------------------
# Standalone functions (backward compatibility)
# ---------------------------------------------------------------------------

def classify_document(file_bytes: bytes, filename: str) -> dict:
    """
    Classifica o tipo de documento financeiro utilizando Claude Sonnet.

    Args:
        file_bytes: Conteudo binario do arquivo.
        filename: Nome do arquivo (usado para inferir formato).

    Returns:
        Dict com chaves: tipo, confianca, descricao.
    """
    try:
        client = _get_client()

        prompt = (
            "Voce e um analista de credito estruturado da ZYN Capital. "
            "Analise o documento fornecido e classifique-o em UMA das categorias abaixo.\n\n"
            "Categorias possiveis:\n"
            "- balanco: Balanco Patrimonial (isolado)\n"
            "- dre: DRE isolada\n"
            "- balancete: Balancete contabil\n"
            "- demonstracoes_financeiras: Demonstracoes Financeiras completas (BP + DRE + DMPL + DFC + Notas)\n"
            "- matricula: Matricula de imovel rural ou urbano\n"
            "- contrato: Contrato (emprestimo, arrendamento, compra e venda, etc.)\n"
            "- certidao: Certidao (negativa de debitos, protestos, distribuicao, etc.)\n"
            "- ccir_car: CCIR ou CAR (Cadastro Ambiental Rural)\n"
            "- laudo_avaliacao: Laudo ou Relatorio Tecnico de Avaliacao de imoveis/ativos\n"
            "- irpf: Declaracao de Imposto de Renda Pessoa Fisica (DIRPF)\n"
            "- faturamento: Relatorio de faturamento, receita ou analise de receita\n"
            "- planejamento: Planejamento de producao, orcamento ou projecao\n"
            "- cnpj: Cartao CNPJ, comprovante de inscricao ou consulta CNPJ\n"
            "- alteracao_contratual: Alteracao contratual, aditivo ou consolidacao societaria\n"
            "- outro: Documento que nao se encaixe em nenhuma categoria acima\n\n"
            "Se o documento contiver BP + DRE juntos, classifique como 'demonstracoes_financeiras'.\n\n"
            "Responda APENAS com JSON, sem texto adicional:\n"
            '{"tipo": "<categoria>", "confianca": <0.0 a 1.0>, "descricao": "<breve descricao>"}'
        )

        content = _build_content_blocks(file_bytes, filename, prompt)

        response_text = _call_api(client, content, max_tokens=4096)

        result = _parse_json_response(response_text)

        # Validacao basica
        if result.get("tipo") not in TIPOS_DOCUMENTO:
            result["tipo"] = "outro"
        if "confianca" not in result:
            result["confianca"] = 0.0
        if "descricao" not in result:
            result["descricao"] = ""

        return result

    except Exception as e:
        logger.exception("Erro ao classificar documento '%s'", filename)
        return {"tipo": "outro", "confianca": 0.0, "descricao": "", "error": str(e)}


def extract_data(file_bytes: bytes, filename: str, tipo_documento: str) -> dict:
    """
    Extrai dados estruturados de um documento financeiro utilizando Claude Sonnet.

    Args:
        file_bytes: Conteudo binario do arquivo.
        filename: Nome do arquivo.
        tipo_documento: Tipo do documento (balanco, dre, matricula, contrato, etc.).

    Returns:
        Dict com os dados extraidos conforme schema do tipo de documento.
    """
    try:
        client = _get_client()

        prompt = _get_extraction_prompt(tipo_documento)
        content = _build_content_blocks(file_bytes, filename, prompt)

        response_text = _call_api(client, content, max_tokens=4096)

        result = _parse_json_response(response_text)
        result["_tipo_documento"] = tipo_documento
        return result

    except Exception as e:
        logger.exception("Erro ao extrair dados do documento '%s'", filename)
        return {"error": str(e), "_tipo_documento": tipo_documento}


# ---------------------------------------------------------------------------
# Unified process_file (1 API call instead of 2)
# ---------------------------------------------------------------------------

def process_file(file_bytes: bytes, filename: str) -> dict:
    """
    Classifica o documento e extrai os dados em UMA unica chamada API.

    Args:
        file_bytes: Conteudo binario do arquivo.
        filename: Nome do arquivo.

    Returns:
        Dict com chaves:
            - classificacao: resultado da classificacao
            - dados: dados extraidos
    """
    size_kb = len(file_bytes) / 1024
    logger.info("[EXTRACT] Iniciando: %s (%.1f KB)", filename, size_kb)

    # Check cache first
    file_hash = _get_file_hash(file_bytes)
    cached = _get_cached(file_hash)
    if cached is not None:
        logger.info("[EXTRACT] Cache hit: %s", filename)
        cached["_from_cache"] = True
        return cached

    try:
        logger.info("[EXTRACT] Preparando conteúdo: %s", filename)
        client = _get_client()

        # Unified prompt: classify + extract in 1 call
        prompt = _build_unified_prompt(TIPOS_DOCUMENTO)
        content = _build_content_blocks(file_bytes, filename, prompt)

        # Log content size for debugging
        total_content_len = sum(len(str(b)) for b in content)
        logger.info("[EXTRACT] Chamando API: %s (content ~%d chars)", filename, total_content_len)

        response_text = _call_api(client, content, max_tokens=4096)
        result = _parse_json_response(response_text)

        # Parse unified response
        if "classificacao" in result and "dados" in result:
            classificacao = result["classificacao"]
            dados = result["dados"]
        else:
            # Fallback: se o modelo retornou formato inesperado, tenta interpretar
            classificacao = {
                "tipo": result.get("tipo", "outro"),
                "confianca": result.get("confianca", 0.0),
                "descricao": result.get("descricao", ""),
            }
            dados = {k: v for k, v in result.items() if k not in ("tipo", "confianca", "descricao", "error")}

        # Validacao basica da classificacao
        if classificacao.get("tipo") not in TIPOS_DOCUMENTO:
            classificacao["tipo"] = "outro"
        if "confianca" not in classificacao:
            classificacao["confianca"] = 0.0
        if "descricao" not in classificacao:
            classificacao["descricao"] = ""

        dados["_tipo_documento"] = classificacao["tipo"]

        logger.info(
            "Documento classificado como '%s' (confianca: %.2f)",
            classificacao["tipo"],
            classificacao.get("confianca", 0.0),
        )

        output = {
            "classificacao": classificacao,
            "dados": dados,
        }

        # Save to cache
        _save_cache(file_hash, output)

        return output

    except Exception as e:
        logger.exception("Erro ao processar documento '%s'", filename)
        classificacao = {"tipo": "outro", "confianca": 0.0, "descricao": "", "error": str(e)}
        return {"classificacao": classificacao, "dados": {"error": str(e)}}


# ---------------------------------------------------------------------------
# Parallel processing
# ---------------------------------------------------------------------------

def process_files_parallel(
    files: list[tuple[bytes, str]],  # (file_bytes, filename)
    max_workers: int = 3,
    progress_callback=None,  # called with (filename, idx, total, result)
) -> dict[str, dict]:
    """Process multiple files in parallel. Returns {filename: result}."""
    total = len(files)
    results: dict[str, dict] = {}

    if total == 0:
        return results

    logger.info("[PARALLEL] Iniciando %d arquivos (max_workers=%d)", total, max_workers)

    _completed = {"count": 0}

    def _process_one(args: tuple[int, bytes, str]) -> tuple[str, dict]:
        idx, file_bytes, filename = args
        logger.info("[PARALLEL] [%d/%d] Processando: %s", idx, total, filename)
        result = process_file(file_bytes, filename)
        _completed["count"] += 1
        logger.info("[PARALLEL] [%d/%d] Concluído: %s (total done: %d)", idx, total, filename, _completed["count"])
        if progress_callback:
            try:
                progress_callback(filename, _completed["count"], total, result)
            except Exception as cb_err:
                logger.warning("Erro no progress_callback para '%s': %s", filename, cb_err)
        return filename, result

    work_items = [(i + 1, fb, fn) for i, (fb, fn) in enumerate(files)]

    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_file = {
            executor.submit(_process_one, item): item[2]
            for item in work_items
        }

        for future in concurrent.futures.as_completed(future_to_file):
            filename = future_to_file[future]
            try:
                fname, result = future.result(timeout=180)  # 3 min max per file
                results[fname] = result
            except concurrent.futures.TimeoutError:
                logger.error("[PARALLEL] TIMEOUT ao processar '%s' (>180s)", filename)
                results[filename] = {
                    "classificacao": {"tipo": "outro", "confianca": 0.0, "descricao": "", "error": "Timeout (>180s)"},
                    "dados": {"error": f"Timeout ao processar {filename} (>180s)"},
                }
            except Exception as e:
                logger.exception("[PARALLEL] Erro fatal ao processar '%s'", filename)
                results[filename] = {
                    "classificacao": {"tipo": "outro", "confianca": 0.0, "descricao": "", "error": str(e)},
                    "dados": {"error": str(e)},
                }

    logger.info("[PARALLEL] Concluído: %d/%d arquivos processados", len(results), total)
    return results
