"""
ZYN Capital — Gerador de Teaser (.pptx) baseado em template oficial.
v2.1 — Deep extraction com suporte a formato direto e wrapper.

Abre o template Teaser_Template_ZYN.pptx (5 slides, 10x5.62 in, 16:9) e
substitui os placeholders pelo conteudo real da operacao, preservando toda
a formatacao visual (fontes, cores, tamanhos, posicoes).

Template:
    Slide 1 — Cover
    Slide 2 — Resumo Executivo & Termos Indicativos (KPI cards + tabela)
    Slide 3 — Overview da Empresa (info + KPI cards)
    Slide 4 — Estrutura da Operacao & Indicadores Financeiros (flow + tabelas)
    Slide 5 — Garantias & Soundness (6 cards + disclaimer)
"""

from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Constantes
# ---------------------------------------------------------------------------
TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "Teaser_Template_ZYN.pptx"

MESES_PT = {
    1: "Janeiro", 2: "Fevereiro", 3: "Marco", 4: "Abril",
    5: "Maio", 6: "Junho", 7: "Julho", 8: "Agosto",
    9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro",
}


# ---------------------------------------------------------------------------
# Helpers — Formatacao BR (mantidos para compatibilidade)
# ---------------------------------------------------------------------------
def _fmt_brl(value, suffix: str = "") -> str:
    """Formata valor em Reais no padrao brasileiro."""
    try:
        v = float(value)
    except (TypeError, ValueError):
        return str(value) if value else "—"
    if abs(v) >= 1_000_000_000:
        formatted = f"R$ {v / 1_000_000_000:,.2f} Bi".replace(",", "X").replace(".", ",").replace("X", ".")
    elif abs(v) >= 1_000_000:
        formatted = f"R$ {v / 1_000_000:,.1f} MM".replace(",", "X").replace(".", ",").replace("X", ".")
    elif abs(v) >= 1_000:
        formatted = f"R$ {v / 1_000:,.0f} mil".replace(",", "X").replace(".", ",").replace("X", ".")
    else:
        formatted = f"R$ {v:,.0f}".replace(",", "X").replace(".", ",").replace("X", ".")
    return formatted + (f" {suffix}" if suffix else "")


def _fmt_pct(value, decimals: int = 1) -> str:
    """Formata percentual no padrao brasileiro."""
    try:
        v = float(value)
        if 0 < abs(v) <= 1:
            v = v * 100
        return f"{v:,.{decimals}f}%".replace(",", "X").replace(".", ",").replace("X", ".")
    except (TypeError, ValueError):
        return str(value) if value else "—"


def _fmt_mult(value) -> str:
    """Formata multiplo (ex: 2,34x)."""
    try:
        v = float(value)
        return f"{v:,.2f}x".replace(",", "X").replace(".", ",").replace("X", ".")
    except (TypeError, ValueError):
        return str(value) if value else "—"


def _safe_get(d, *keys, default="—"):
    """Navegacao segura em dicts aninhados. Retorna default se qualquer nivel nao for dict."""
    current = d
    for k in keys:
        if isinstance(current, dict):
            current = current.get(k, default)
        else:
            return default
    return current if current not in (None, "", 0, 0.0) else default


def _safe_dict(value, default=None):
    """Garante que o valor é um dict; retorna default={} se nao for."""
    return value if isinstance(value, dict) else (default or {})


def _current_date_pt() -> str:
    """Retorna 'Marco 2026' no formato portugues."""
    now = datetime.now()
    return f"{MESES_PT.get(now.month, '')} {now.year}"


# ---------------------------------------------------------------------------
# Core: substituicao de texto preservando formatacao
# ---------------------------------------------------------------------------
def _replace_text_in_runs(text_frame, old_text: str, new_text: str) -> bool:
    """
    Substitui texto em runs de um text_frame, preservando formatacao.
    Tenta primeiro run-a-run; se nao encontrar, tenta concatenar runs
    de cada paragrafo para lidar com texto dividido entre runs.
    """
    replaced = False
    for para in text_frame.paragraphs:
        # Tentativa 1: substituicao direta em cada run
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                replaced = True
        # Tentativa 2: texto dividido entre runs adjacentes
        if not replaced:
            full_text = "".join(r.text for r in para.runs)
            if old_text in full_text:
                new_full = full_text.replace(old_text, new_text)
                # Coloca todo o texto no primeiro run, limpa os demais
                if para.runs:
                    para.runs[0].text = new_full
                    for run in para.runs[1:]:
                        run.text = ""
                    replaced = True
    return replaced


def _replace_in_shape(shape, old_text: str, new_text: str) -> bool:
    """Substitui texto em shape (text_frame), preservando formatacao."""
    if not shape.has_text_frame:
        return False
    return _replace_text_in_runs(shape.text_frame, old_text, new_text)


def _replace_in_table_cell(table, row: int, col: int, new_text: str):
    """Substitui todo o conteudo de uma celula de tabela, preservando formatacao do primeiro run."""
    try:
        cell = table.cell(row, col)
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].text = str(new_text)
            for run in cell.text_frame.paragraphs[0].runs[1:]:
                run.text = ""
        else:
            cell.text = str(new_text)
    except (IndexError, AttributeError):
        pass


def _replace_in_table(table, old_text: str, new_text: str) -> bool:
    """Busca e substitui texto em qualquer celula da tabela."""
    replaced = False
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            try:
                if _replace_text_in_runs(cell.text_frame, old_text, new_text):
                    replaced = True
            except Exception:
                pass
    return replaced


def _replace_on_slide(slide, old_text: str, new_text: str) -> bool:
    """Substitui texto em todas as shapes (incluindo tabelas) de um slide."""
    replaced = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            if _replace_in_shape(shape, old_text, new_text):
                replaced = True
        if shape.has_table:
            if _replace_in_table(shape.table, old_text, new_text):
                replaced = True
    return replaced


# ---------------------------------------------------------------------------
# Helpers — Deep extraction from analise dict
# ---------------------------------------------------------------------------
def _deep_get(analise: dict, *paths, default="—"):
    """
    Busca um valor em multiplos caminhos possiveis no dict de analise.
    O Opus retorna dados em analise.analise.X ou analise.X dependendo do contexto.
    Tenta todos os caminhos ate achar um valor valido.
    """
    for path in paths:
        val = _safe_get(analise, *path, default=None)
        if val not in (None, "", 0, 0.0, "—", []):
            return val
    return default


def _flatten_grupo_economico(ge) -> str:
    """Converte dict grupo_economico em texto profissional para investidores."""
    if isinstance(ge, str):
        return ge
    if not isinstance(ge, dict):
        return "—"
    parts = []
    # Sócios PF com nomes limpos
    socios_pf = ge.get("socios_pf", [])
    if isinstance(socios_pf, list) and socios_pf:
        nomes = []
        for s in socios_pf:
            if isinstance(s, str):
                # Extrai nome limpo: "MARCO AURELIO FUENTES HOLLATZ (CPF ...) - 50%..."
                nome = s.split("(")[0].strip().title()
                pct = ""
                if "%" in s:
                    import re
                    m = re.search(r'(\d+%)', s)
                    if m:
                        pct = f" ({m.group(1)})"
                nomes.append(f"{nome}{pct}")
            elif isinstance(s, dict):
                nomes.append(s.get("nome", str(s)))
        if nomes:
            parts.append(", ".join(nomes))
    # Holding
    holding = ge.get("holding", "")
    if isinstance(holding, str) and holding:
        nome_holding = holding.split("(")[0].strip()
        if nome_holding:
            parts.append(f"via {nome_holding}")
    # Usufrutuários / fundadores
    usuf = ge.get("usufrutuarios", [])
    if isinstance(usuf, list) and usuf:
        for u in usuf:
            if isinstance(u, str) and "fundador" in u.lower():
                nome_f = u.split("(")[0].strip().title()
                parts.append(f"Fundador(a): {nome_f}")
                break
    return " | ".join(parts) if parts else "—"


def _flatten_historico(hist) -> str:
    """Converte dict historico em narrativa profissional."""
    if isinstance(hist, str):
        return hist
    if not isinstance(hist, dict):
        return "—"
    parts = []
    anos = hist.get("anos_operacao")
    fund = hist.get("fundacao_grupo") or hist.get("fundacao")
    if anos:
        parts.append(f"{anos} anos de atuação")
    elif fund:
        parts.append(f"Fundado em {fund}")
    entregues = hist.get("empreendimentos_entregues") or hist.get("projetos_entregues")
    if entregues:
        parts.append(f"{entregues} empreendimentos entregues")
    unidades = hist.get("unidades_construidas") or hist.get("unidades_historicas")
    if unidades:
        parts.append(f"+{unidades:,} unidades".replace(",", "."))
    m2 = hist.get("metros_quadrados_construidos")
    if m2:
        parts.append(f"+{m2:,} m² construídos".replace(",", "."))
    pontualidade = hist.get("taxa_pontualidade")
    if pontualidade:
        parts.append(f"{pontualidade} de pontualidade")
    satisfacao = hist.get("satisfacao_clientes")
    if satisfacao:
        parts.append(f"{satisfacao} satisfação dos clientes")
    certs = hist.get("certificacoes", [])
    if isinstance(certs, list) and certs:
        cert_names = [c.split("(")[0].strip() if isinstance(c, str) else str(c) for c in certs[:3]]
        parts.append(f"Certificações: {', '.join(cert_names)}")
    obras = hist.get("obras_em_andamento")
    vgv_obras = hist.get("vgv_obras_andamento")
    if obras and vgv_obras:
        parts.append(f"{obras} obras em andamento (VGV {_fmt_brl(vgv_obras)})")
    elif obras:
        parts.append(f"{obras} obras em andamento")
    return ". ".join(parts) + "." if parts else "—"


def _flatten_capacidade(cap) -> str:
    """Converte dict capacidade em texto legível."""
    if isinstance(cap, str):
        return cap
    if not isinstance(cap, dict):
        return "—"
    parts = []
    obras_sim = cap.get("obras_simultaneas_atuais")
    if obras_sim:
        parts.append(f"{obras_sim} obras simultâneas")
    vgv_and = cap.get("vgv_total_em_andamento")
    if vgv_and:
        parts.append(f"VGV em andamento: {_fmt_brl(vgv_and)}")
    un_and = cap.get("unidades_totais_em_andamento")
    if un_and:
        parts.append(f"{un_and} unidades em andamento")
    pipeline = cap.get("pipeline_futuro_vgv")
    if pipeline:
        parts.append(f"Pipeline futuro: {_fmt_brl(pipeline)}")
    conc = cap.get("obras_concluidas_historicas")
    if conc:
        parts.append(f"{conc} obras concluídas")
    m2_hist = cap.get("m2_historicos")
    if m2_hist:
        parts.append(f"{m2_hist:,} m² históricos".replace(",", "."))
    # Agro
    hectares = cap.get("hectares") or cap.get("area_total_ha")
    if hectares:
        parts.append(f"{hectares:,} ha".replace(",", "."))
    equip = cap.get("equipamentos")
    if equip:
        parts.append(f"{equip} equipamentos")
    return " | ".join(parts) if parts else "—"


def _flatten_historico_produtivo(hp) -> str:
    """Extrai dados do empreendimento principal do historico produtivo."""
    if isinstance(hp, str):
        return hp
    if not isinstance(hp, dict):
        return "—"
    # Busca o primeiro empreendimento detalhado
    for key, val in hp.items():
        if isinstance(val, dict) and "total_unidades" in val:
            parts = []
            un = val.get("total_unidades")
            vgv = val.get("vgv_total")
            vendidas = val.get("vendidas")
            vso = val.get("vso_atual_pct")
            avanco = val.get("avanco_fisico_pct")
            ticket = val.get("ticket_medio")
            if un:
                parts.append(f"{un} unidades")
            if vgv:
                parts.append(f"VGV {_fmt_brl(vgv)}")
            if vendidas and vso:
                parts.append(f"{vendidas} vendidas ({vso:.0f}% VSO)")
            if ticket:
                parts.append(f"Ticket médio {_fmt_brl(ticket)}")
            if avanco:
                parts.append(f"Avanço físico {avanco:.0f}%")
            return " | ".join(parts) if parts else "—"
    return "—"


def _flatten_any(val) -> str:
    """Converte qualquer valor para string legível. Último recurso."""
    if isinstance(val, str):
        return val
    if isinstance(val, (int, float)):
        return str(val)
    if isinstance(val, list):
        items = []
        for item in val[:5]:
            if isinstance(item, str):
                items.append(item)
            elif isinstance(item, dict):
                items.append(str(item.get("nome", item.get("descricao", str(item)))))
            else:
                items.append(str(item))
        return "; ".join(items)
    if isinstance(val, dict):
        # Tenta serializar campos-chave
        parts = []
        for k, v in list(val.items())[:6]:
            if isinstance(v, (str, int, float)) and v:
                parts.append(f"{k}: {v}")
        return " | ".join(parts) if parts else str(val)
    return str(val) if val else "—"


def _extract_company_data(analise: dict, parametros: dict) -> dict:
    """Extrai todos os dados da empresa de todas as secoes da analise.

    Suporta dois formatos:
    1. Dict direto do Opus: {kpis: {}, tomador: {}, rating_final: {}, ...}
    2. Dict salvo com wrapper: {analise: {kpis: {}, ...}, operacao: {...}}

    Todos os valores retornados sao strings ou numeros — nunca dicts/lists crus.
    """
    # Detecta se tem camada wrapper ou se é direto
    if isinstance(analise.get("analise"), dict) and "kpis" in analise.get("analise", {}):
        a = _safe_dict(analise.get("analise"))
    elif "kpis" in analise or "tomador" in analise or "rating_final" in analise:
        a = _safe_dict(analise)
    else:
        a = {}

    tom_a = _safe_dict(a.get("tomador"))
    tom_root = _safe_dict(analise.get("tomador"))
    op_root = _safe_dict(analise.get("operacao")) if isinstance(analise.get("operacao"), dict) else {}
    op_a = _safe_dict(a.get("operacao"))
    kpis_a = _safe_dict(a.get("kpis"))
    kpis_root = _safe_dict(analise.get("kpis"))
    cap = _safe_dict(a.get("capital"))
    cap_struct = _safe_dict(cap.get("estrutura_capital"))
    cap_ind = _safe_dict(cap.get("indicadores"))
    cap_end = _safe_dict(cap.get("endividamento"))
    pat = _safe_dict(a.get("patrimonio"))
    pat_ativos = _safe_dict(pat.get("ativos_reais"))
    prod = _safe_dict(a.get("producao"))
    pag = _safe_dict(a.get("pagamento"))
    rating_a = _safe_dict(a.get("rating_final")) if isinstance(a.get("rating_final"), dict) else {}
    rating_root = _safe_dict(analise.get("rating_final")) if isinstance(analise.get("rating_final"), dict) else {}
    rating = rating_a if rating_a else rating_root

    # --- Socios: flatten grupo_economico ---
    socios_raw = tom_a.get("grupo_economico") or tom_a.get("socios") or tom_root.get("socios") or tom_root.get("grupo_economico")
    if isinstance(socios_raw, dict):
        socios = _flatten_grupo_economico(socios_raw)
    elif isinstance(socios_raw, list):
        socios = "; ".join(str(s).split("(")[0].strip().title() if isinstance(s, str) else str(s) for s in socios_raw[:4])
    elif isinstance(socios_raw, str):
        socios = socios_raw
    else:
        socios = "—"

    # --- Descricao: flatten historico dict ---
    desc_raw = tom_a.get("historico") or tom_root.get("descricao") or tom_root.get("historico")
    if isinstance(desc_raw, dict):
        descricao = _flatten_historico(desc_raw)
    elif isinstance(desc_raw, str):
        descricao = desc_raw
    else:
        # Fallback: build from capacidade + producao
        cap_raw = prod.get("capacidade")
        if isinstance(cap_raw, dict):
            descricao = _flatten_capacidade(cap_raw)
        elif isinstance(cap_raw, str):
            descricao = cap_raw
        else:
            descricao = "—"

    # --- Fundacao: extract from historico dict ---
    fundacao = tom_a.get("fundacao") or tom_root.get("fundacao")
    if not fundacao:
        hist_raw = tom_a.get("historico")
        if isinstance(hist_raw, dict):
            fundacao = hist_raw.get("fundacao_grupo") or hist_raw.get("fundacao") or ""
        elif isinstance(hist_raw, str):
            fundacao = _extract_year_from_text(hist_raw)
        if not fundacao:
            hp_raw = prod.get("historico_produtivo")
            if isinstance(hp_raw, str):
                fundacao = _extract_year_from_text(hp_raw)
    fundacao = str(fundacao) if fundacao else "—"

    # --- Capacidade: flatten ---
    cap_raw = tom_a.get("capacidade") or tom_root.get("capacidade") or prod.get("capacidade")
    if isinstance(cap_raw, dict):
        capacidade = _flatten_capacidade(cap_raw)
    elif isinstance(cap_raw, str):
        capacidade = cap_raw
    else:
        capacidade = "—"

    # --- Unidades: extract from capacidade or producao ---
    un_raw = tom_a.get("unidades") or tom_root.get("unidades")
    if not un_raw:
        cap_d = prod.get("capacidade") if isinstance(prod.get("capacidade"), dict) else {}
        un_and = cap_d.get("unidades_totais_em_andamento")
        obras = cap_d.get("obras_simultaneas_atuais")
        if un_and and obras:
            un_raw = f"{un_and} un. em {obras} obras"
        elif un_and:
            un_raw = f"{un_and} unidades"
        elif isinstance(prod.get("capacidade"), str):
            un_raw = _extract_number_near_keyword(prod["capacidade"], "unidade")
    unidades = str(un_raw) if un_raw else "—"

    # --- Colaboradores ---
    colab_raw = tom_a.get("colaboradores") or tom_root.get("colaboradores")
    if not colab_raw:
        ge = tom_a.get("grupo_economico")
        if isinstance(ge, dict):
            colab_raw = ge.get("total_membros_grupo")
            if colab_raw:
                colab_raw = f"{colab_raw} membros no grupo"
        if not colab_raw:
            hist_str = tom_a.get("historico", "")
            cap_str = prod.get("capacidade", "")
            if isinstance(hist_str, str) and isinstance(cap_str, str):
                colab_raw = _extract_number_near_keyword(hist_str + " " + cap_str, "colaborador")
    colaboradores = str(colab_raw) if colab_raw else "—"

    # --- Clientes ---
    clientes_raw = tom_a.get("principais_clientes") or tom_root.get("principais_clientes")
    if not clientes_raw:
        prod_analise = prod.get("analise", "")
        if isinstance(prod_analise, str):
            clientes_raw = _extract_clients_from_text(prod_analise)
    if isinstance(clientes_raw, list):
        clientes_raw = ", ".join(str(c) for c in clientes_raw[:5])
    clientes = str(clientes_raw) if clientes_raw else "—"

    # --- Amortizacao: flatten from operacao.estrutura ---
    amort_raw = parametros.get("amortizacao") or op_root.get("amortizacao")
    if not amort_raw:
        est = op_a.get("estrutura")
        if isinstance(est, dict):
            amort_raw = est.get("amortizacao", "")
        elif isinstance(est, str):
            amort_raw = est
    amortizacao = str(amort_raw) if amort_raw else "—"

    # --- Carencia ---
    carencia_raw = parametros.get("carencia")
    if not carencia_raw:
        est = op_a.get("estrutura")
        if isinstance(est, dict):
            carencia_raw = est.get("carencia", "")
        if not carencia_raw:
            est_str = est if isinstance(est, str) else ""
            prazo_str = op_a.get("prazo", "") if isinstance(op_a.get("prazo"), str) else ""
            carencia_raw = _extract_from_text(est_str + " " + prazo_str, r"(?:[Cc]ar[eê]ncia\s*(?:de\s*)?)(\d+\s*meses?)")
    carencia = str(carencia_raw) if carencia_raw else "—"

    # --- Finalidade: prioriza lastro/destinacao sobre analise de risco ---
    final_raw = parametros.get("finalidade") or parametros.get("destinacao_texto")
    if not final_raw:
        est = op_a.get("estrutura")
        if isinstance(est, dict):
            lastro = est.get("lastro", "")
            if lastro:
                final_raw = str(lastro)
    if not final_raw:
        op_analise = op_a.get("analise")
        if isinstance(op_analise, str):
            # Pega primeira frase relevante, evita texto de risco
            sentences = op_analise.split(". ")
            for s in sentences:
                s_lower = s.lower()
                if any(kw in s_lower for kw in ["lastro", "destin", "financ", "capta", "emissão", "operação de"]):
                    final_raw = s.strip()[:250]
                    break
            if not final_raw:
                final_raw = sentences[0][:200] if sentences else ""
    finalidade = str(final_raw) if final_raw else "—"

    # --- PL: busca em capital.estrutura_capital ---
    pl = pat_ativos.get("patrimonio_liquido_2025") or pat_ativos.get("patrimonio_liquido") or kpis_a.get("pl") or 0
    if not pl:
        for entity_key, entity_val in cap_struct.items():
            if isinstance(entity_val, dict):
                pl_candidate = entity_val.get("patrimonio_liquido")
                if pl_candidate and isinstance(pl_candidate, (int, float)) and pl_candidate > 0:
                    pl = pl_candidate
                    break

    # --- Empreendimento principal (para teaser imobiliário) ---
    hp = prod.get("historico_produtivo")
    empreendimento_str = ""
    if isinstance(hp, dict):
        empreendimento_str = _flatten_historico_produtivo(hp)

    # --- Garantias detalhadas ---
    gar_det = pat.get("garantias_detalhadas", [])
    if not isinstance(gar_det, list):
        gar_det = []

    return {
        # Company info
        "nome": parametros.get("tomador") or tom_a.get("razao_social") or tom_root.get("nome") or tom_root.get("razao_social") or "—",
        "cnpj": parametros.get("cnpj") or tom_a.get("cnpj") or tom_root.get("cnpj") or "—",
        "fundacao": fundacao,
        "sede": parametros.get("localidade") or tom_a.get("sede") or tom_root.get("sede") or _extract_city_from_text(descricao if isinstance(descricao, str) else "") or _extract_city_from_text(str(rating.get("justificativa", ""))) or "—",
        "segmento": parametros.get("setor") or tom_a.get("setor") or tom_root.get("segmento") or tom_root.get("setor") or "—",
        "socios": socios,
        "descricao": descricao,
        "colaboradores": colaboradores,
        "capacidade": capacidade,
        "unidades": unidades,
        "clientes": clientes,
        "empreendimento": empreendimento_str,

        # Financials
        "receita": kpis_a.get("receita_liquida") or kpis_root.get("receita_liquida") or 0,
        "ebitda": kpis_a.get("ebitda") or kpis_root.get("ebitda") or 0,
        "margem_ebitda": kpis_a.get("margem_ebitda") or kpis_root.get("margem_ebitda") or 0,
        "divida_liquida": cap_end.get("divida_liquida_2025") or cap_end.get("divida_liquida") or kpis_a.get("divida_liquida") or 0,
        "div_liq_ebitda": cap_ind.get("divida_liquida_ebitda") or kpis_a.get("divida_liquida_ebitda") or kpis_root.get("divida_liquida_ebitda") or 0,
        "dscr": pag.get("dscr") or kpis_a.get("dscr") or kpis_root.get("dscr") or 0,
        "ltv": kpis_a.get("ltv") or kpis_root.get("ltv") or pat.get("ltv") or 0,
        "pl": pl,
        "liquidez": cap_ind.get("liquidez_corrente") or 0,

        # Operation
        "instrumento": parametros.get("tipo_operacao") or op_root.get("tipo_operacao") or op_a.get("instrumento") or op_root.get("instrumento") or "—",
        "volume": parametros.get("volume") or op_root.get("volume") or op_a.get("volume") or 0,
        "taxa": parametros.get("taxa") or op_root.get("taxa") or op_a.get("taxa") or "—",
        "prazo": parametros.get("prazo_meses") or op_root.get("prazo_meses") or op_a.get("prazo") or "—",
        "amortizacao": amortizacao,
        "carencia": carencia,
        "garantias_text": parametros.get("garantias_text") or op_root.get("garantias_text") or "—",
        "finalidade": finalidade,

        # Guarantees (detailed)
        "garantias_detalhadas": gar_det,

        # Rating
        "rating": rating.get("nota") or op_root.get("rating") or "—",
        "parecer": rating.get("parecer") or op_root.get("parecer") or "—",
        "justificativa": rating.get("justificativa") or "—",
    }


def _extract_from_text(text: str, pattern: str) -> str:
    """Extrai valor de texto usando regex. Retorna '' se nao encontrar."""
    import re
    if not isinstance(text, str):
        return ""
    match = re.search(pattern, text)
    return match.group(1) if match else ""


def _extract_year_from_text(text: str) -> str:
    """Extrai ano de fundacao de texto tipo 'Fundada em 1978' ou 'desde 1985'."""
    import re
    if not isinstance(text, str) or not text:
        return ""
    for pattern in [r'[Ff]undad[ao]\s+em\s+(\d{4})', r'desde\s+(\d{4})', r'criada?\s+em\s+(\d{4})', r'(\d{4})\s+em\s+\w']:
        m = re.search(pattern, text)
        if m:
            year = m.group(1)
            if 1900 <= int(year) <= 2030:
                return year
    return ""


def _extract_city_from_text(text: str) -> str:
    """Extrai cidade/UF de texto tipo 'em Chapeco-SC' ou 'sediada em Sao Paulo/SP'."""
    import re
    if not isinstance(text, str) or not text:
        return ""
    # Tenta padroes comuns: Cidade-UF, Cidade/UF, Cidade (UF)
    for pattern in [r'em\s+([A-Z][\w\u00C0-\u00FF]+(?:\s+[\w\u00C0-\u00FF]+)*\s*[-/]\s*[A-Z]{2})',
                    r'sediada?\s+em\s+([\w\u00C0-\u00FF]+(?:\s+[\w\u00C0-\u00FF]+)*)',
                    r'sede\s+(?:em\s+)?([\w\u00C0-\u00FF]+[-/][A-Z]{2})']:
        m = re.search(pattern, text)
        if m:
            return m.group(1).strip()
    return ""


def _extract_number_near_keyword(text: str, keyword: str) -> str:
    """Extrai numero proximo a uma keyword. Ex: '~3.000 equipamentos' -> '3.000'."""
    import re
    if not isinstance(text, str) or not text or not keyword:
        return ""
    # Busca: ~N.NNN keyword ou keyword N.NNN
    pattern = rf'~?(\d[\d.]*)\s*{keyword}'
    m = re.search(pattern, text, re.IGNORECASE)
    if m:
        return m.group(1)
    pattern = rf'{keyword}\w*\s+~?(\d[\d.]*)'
    m = re.search(pattern, text, re.IGNORECASE)
    if m:
        return m.group(1)
    return ""


def _extract_clients_from_text(text: str) -> str:
    """Extrai nomes de clientes de texto tipo '(Ambev, Heineken, JBS, Suzano)'."""
    import re
    if not isinstance(text, str) or not text:
        return ""
    # Busca lista entre parenteses com nomes proprios
    m = re.search(r'\(([A-Z][\w]+(?:,\s*[A-Z][\w]+){2,})\)', text)
    if m:
        return m.group(1)
    return ""


# ---------------------------------------------------------------------------
# Slide 1 — Cover
# ---------------------------------------------------------------------------
def _fill_cover(slide, analise: dict, parametros: dict):
    data = _extract_company_data(analise, parametros)
    data_pt = _current_date_pt()

    _replace_on_slide(slide, "[TIPO DE INSTRUMENTO]", str(data["instrumento"]).upper())
    _replace_on_slide(slide, "[Nome da Operação / Tomador]", str(data["nome"]))
    _replace_on_slide(slide, "Março 2026", data_pt)
    _replace_on_slide(slide, "Marco 2026", data_pt)


# ---------------------------------------------------------------------------
# Slide 2 — Resumo Executivo & Termos Indicativos
# ---------------------------------------------------------------------------
def _fill_resumo(slide, analise: dict, parametros: dict):
    data = _extract_company_data(analise, parametros)
    data_pt = _current_date_pt()
    _replace_on_slide(slide, "Março 2026", data_pt)
    _replace_on_slide(slide, "Marco 2026", data_pt)

    # --- KPI cards ---
    volume_str = _fmt_brl(data["volume"]) if data["volume"] else "—"
    taxa_str = str(data["taxa"]) if data["taxa"] != "—" else "—"
    prazo_str = f"{data['prazo']}" if data["prazo"] != "—" else "—"
    if prazo_str != "—" and "mes" not in str(prazo_str).lower():
        prazo_str = f"{prazo_str} meses"
    ltv_str = _fmt_pct(data["ltv"]) if data["ltv"] else "—"

    _replace_on_slide(slide, "R$ [XX] MM", volume_str)
    _replace_on_slide(slide, "CDI + [X]%", taxa_str)
    _replace_on_slide(slide, "[XX] meses", prazo_str)
    _replace_on_slide(slide, "[XX]%", ltv_str)

    # --- Tabela de termos (11 rows x 2 cols, header em row 0) ---
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break
    if table is None:
        return

    termos = [
        str(data["nome"]),
        str(data["instrumento"]),
        volume_str,
        taxa_str,
        prazo_str,
        str(data["amortizacao"]),
        str(data["carencia"]),
        str(data["garantias_text"]),
        str(data["finalidade"])[:200],
        str(data["rating"]),
    ]
    for i, val in enumerate(termos):
        _replace_in_table_cell(table, i + 1, 1, val)


# ---------------------------------------------------------------------------
# Slide 3 — Overview da Empresa
# ---------------------------------------------------------------------------
def _fill_overview(slide, analise: dict, parametros: dict):
    data = _extract_company_data(analise, parametros)
    data_pt = _current_date_pt()
    _replace_on_slide(slide, "Março 2026", data_pt)
    _replace_on_slide(slide, "Marco 2026", data_pt)

    # --- Left panel: company info ---
    _replace_on_slide(slide, "[NOME DO GRUPO / EMPRESA]", str(data["nome"]))
    _replace_on_slide(slide, "[Ano]", str(data["fundacao"]))
    _replace_on_slide(slide, "[Cidade/UF]", str(data["sede"]))
    _replace_on_slide(slide, "[Agronegócio / Imobiliário / Industrial / etc.]", str(data["segmento"]))
    _replace_on_slide(slide, "[Nomes e cargos principais]", str(data["socios"])[:300])
    # Descricao: usar narrativa rica, com empreendimento se disponível
    desc = str(data["descricao"])[:400]
    empreend = data.get("empreendimento", "")
    if empreend and empreend != "—":
        desc = f"{desc}\n\nEmpreendimento: {empreend}"
    _replace_on_slide(
        slide,
        "[Breve histórico da empresa, atividades principais, diferenciais competitivos, principais clientes/offtakers, e posicionamento de mercado. 3-4 linhas.]",
        desc[:600],
    )

    # --- Right panel: KPI cards ---
    _replace_on_slide(slide, "R$ [XXX] MM", _fmt_brl(data["receita"]) if data["receita"] else "—")
    _replace_on_slide(slide, "R$ [XX] MM", _fmt_brl(data["ebitda"]) if data["ebitda"] else "—")
    _replace_on_slide(slide, "[XXX]", str(data["colaboradores"]))
    # Capacidade: truncar para caber
    cap_display = str(data["capacidade"])[:80]
    _replace_on_slide(slide, "[XX.XXX ha / m² / ton]", cap_display)
    _replace_on_slide(slide, "[X] unidades em [UFs]", str(data["unidades"]))
    _replace_on_slide(slide, "[Cliente A, B, C]", str(data["clientes"]))


# ---------------------------------------------------------------------------
# Slide 4 — Estrutura da Operacao & Indicadores Financeiros
# ---------------------------------------------------------------------------
def _fill_estrutura(slide, analise: dict, parametros: dict):
    data = _extract_company_data(analise, parametros)
    data_pt = _current_date_pt()
    _replace_on_slide(slide, "Março 2026", data_pt)
    _replace_on_slide(slide, "Marco 2026", data_pt)

    # --- Flow boxes ---
    tomador = str(data["nome"])
    tipo_upper = str(data["instrumento"]).upper()

    # Determine veiculo com base no tipo
    veiculos_map = {
        "CRA": "SECURITIZADORA", "CRI": "SECURITIZADORA",
        "FIDC": "FIDC", "FIAGRO": "FIAGRO",
        "DEBENTURE": "EMISSORA", "NC": "BANCO EMISSOR", "CCB": "BANCO EMISSOR",
        "NC/CCB": "BANCO EMISSOR", "SLB": "OPERADOR S&LB",
    }
    veiculo = veiculos_map.get(tipo_upper, "VEÍCULO")

    _replace_on_slide(slide, "[TOMADOR / CEDENTE]", str(tomador).upper())
    _replace_on_slide(slide, "[SECURITIZADORA / VEÍCULO]", veiculo)
    _replace_on_slide(slide, "[INVESTIDOR]", "INVESTIDOR")

    # --- Destinacao table (shape 19 in template: 4 rows x 3 cols) ---
    destinacao_table = None
    financeiros_table = None
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            ncols = len(tbl.columns)
            if ncols == 3:
                destinacao_table = tbl
            elif ncols == 4:
                financeiros_table = tbl

    # Destinacao: tentamos preencher com dados de parametros ou analise
    if destinacao_table:
        destinacao = parametros.get("destinacao", [])
        if not destinacao:
            finalidade = parametros.get("finalidade", _safe_get(analise, "operacao", "finalidade", default=""))
            volume_raw = parametros.get("volume", _safe_get(analise, "operacao", "volume", default=0))
            if finalidade and finalidade != "—":
                destinacao = [{"destino": str(finalidade), "valor": volume_raw, "pct": "100%"}]

        for i, item in enumerate(destinacao[:3]):  # max 3 data rows
            row_idx = i + 1
            if isinstance(item, str):
                _replace_in_table_cell(destinacao_table, row_idx, 0, item)
                continue
            if not isinstance(item, dict):
                continue
            _replace_in_table_cell(destinacao_table, row_idx, 0, str(item.get("destino", "—")))
            val = item.get("valor", "—")
            try:
                val_mm = f"{float(val) / 1_000_000:,.1f}".replace(",", "X").replace(".", ",").replace("X", ".")
            except (TypeError, ValueError):
                val_mm = str(val)
            _replace_in_table_cell(destinacao_table, row_idx, 1, val_mm)
            _replace_in_table_cell(destinacao_table, row_idx, 2, str(item.get("pct", "—")))

        # Limpa linhas nao usadas
        for i in range(len(destinacao), 3):
            row_idx = i + 1
            _replace_in_table_cell(destinacao_table, row_idx, 0, "—")
            _replace_in_table_cell(destinacao_table, row_idx, 1, "—")
            _replace_in_table_cell(destinacao_table, row_idx, 2, "—")

    # --- Indicadores Financeiros table (9 rows x 4 cols) ---
    if financeiros_table:
        a = _safe_dict(analise.get("analise"))
        kpis = _safe_dict(a.get("kpis")) if a else _safe_dict(analise.get("kpis"))
        cap = _safe_dict(a.get("capital"))
        cap_ind = _safe_dict(cap.get("indicadores"))
        cap_end = _safe_dict(cap.get("endividamento"))
        pat = _safe_dict(a.get("patrimonio"))
        pat_ativos = _safe_dict(pat.get("ativos_reais"))
        pag = _safe_dict(a.get("pagamento"))

        now = datetime.now()
        anos = [str(now.year - 2), str(now.year - 1), str(now.year)]

        # Headers: anos
        for ci, ano in enumerate(anos):
            _replace_in_table_cell(financeiros_table, 0, ci + 1, ano)

        # Build latest year data from all sources
        latest = {
            "receita_liquida": data["receita"],
            "ebitda": data["ebitda"],
            "margem_ebitda": data["margem_ebitda"],
            "divida_liquida": data["divida_liquida"],
            "div_liq_ebitda": data["div_liq_ebitda"],
            "dscr": data["dscr"],
            "ltv": data["ltv"],
            "pl": data["pl"],
        }

        # Try to get historical data
        historico = _safe_dict(analise.get("historico_financeiro"))

        row_map = [
            ("receita_liquida", _fmt_brl),
            ("ebitda", _fmt_brl),
            ("margem_ebitda", _fmt_pct),
            ("divida_liquida", _fmt_brl),
            ("div_liq_ebitda", _fmt_mult),
            ("dscr", _fmt_mult),
            ("ltv", _fmt_pct),
            ("pl", _fmt_brl),
        ]

        for row_idx, (key, fmt) in enumerate(row_map):
            for col_idx, ano in enumerate(anos):
                val = "—"
                # Try historico first
                if historico:
                    val_raw = _safe_get(historico, ano, key, default=None)
                    if val_raw not in (None, "", 0, 0.0, "—"):
                        val = fmt(val_raw)
                # Last column: use latest data
                if val == "—" and col_idx == 2:
                    val_raw = latest.get(key, 0)
                    if val_raw and val_raw not in (0, 0.0, "—"):
                        val = fmt(val_raw)
                _replace_in_table_cell(financeiros_table, row_idx + 1, col_idx + 1, val)


# ---------------------------------------------------------------------------
# Slide 5 — Garantias & Soundness
# ---------------------------------------------------------------------------
def _fill_garantias(slide, analise: dict, parametros: dict):
    data = _extract_company_data(analise, parametros)
    data_pt = _current_date_pt()
    _replace_on_slide(slide, "Março 2026", data_pt)
    _replace_on_slide(slide, "Marco 2026", data_pt)

    # Busca garantias detalhadas de todas as fontes
    garantias_list = data.get("garantias_detalhadas", [])
    if not isinstance(garantias_list, list):
        garantias_list = []
    # Fallback para garantias simples
    if not garantias_list:
        garantias_list = parametros.get("garantias", analise.get("garantias", []))
        if not isinstance(garantias_list, list):
            garantias_list = []

    # Placeholders do template
    ph_alienacao = "[Descrição do imóvel/bem, matrícula, localização, valor de avaliação]"
    ph_cessao = "[Recebíveis cedidos, fluxo, prazo, valor estimado do lastro]"
    ph_aval = "[Avalistas PF/PJ, patrimônio declarado, vínculos com o tomador]"
    ph_fundo = "[X] parcelas equivalentes — constituição [pré/pós] emissão"

    # Classificacao por tipo
    alienacao_descs = []
    cessao_descs = []
    aval_descs = []
    fundo_descs = []

    for gar in garantias_list:
        if isinstance(gar, str):
            tipo_lower = gar.lower()
        elif isinstance(gar, dict):
            tipo_lower = str(gar.get("tipo_garantia", gar.get("tipo", ""))).lower()
        else:
            continue

        # Build description
        if isinstance(gar, dict):
            desc = gar.get("descricao", "")
            valor = gar.get("valor_estimado", 0)
            if valor and valor != 0:
                desc += f" (valor estimado: {_fmt_brl(valor)})"
        else:
            desc = gar

        if any(k in tipo_lower for k in ["alienação", "alienacao", "fiduciária de", "fiduciaria de", "imóv", "imov"]):
            if "cessão" not in tipo_lower and "cessao" not in tipo_lower and "recebív" not in tipo_lower:
                alienacao_descs.append(desc)
        if any(k in tipo_lower for k in ["cessão", "cessao", "recebív", "recebiv", "crédito", "credito"]):
            cessao_descs.append(desc)
        if any(k in tipo_lower for k in ["aval", "fiança", "fianca"]):
            aval_descs.append(desc)
        if any(k in tipo_lower for k in ["fundo", "reserva"]):
            fundo_descs.append(desc)

    if alienacao_descs:
        _replace_on_slide(slide, ph_alienacao, "; ".join(alienacao_descs)[:300])
    if cessao_descs:
        _replace_on_slide(slide, ph_cessao, "; ".join(cessao_descs)[:300])
    if aval_descs:
        _replace_on_slide(slide, ph_aval, "; ".join(aval_descs)[:300])
    if fundo_descs:
        _replace_on_slide(slide, ph_fundo, "; ".join(fundo_descs)[:300])

    # Razao de garantia
    ltv_val = data["ltv"]
    if ltv_val and ltv_val not in (0, 0.0, "—"):
        try:
            ltv_f = float(ltv_val)
            # LTV pode vir como decimal (0.625) ou percentual (62.5)
            if ltv_f > 1:
                ltv_f = ltv_f / 100  # Converte 62.5 -> 0.625
            razao = 1 / ltv_f if 0 < ltv_f < 1 else 1.0
            razao_str = _fmt_mult(razao)
        except (TypeError, ValueError, ZeroDivisionError):
            razao_str = "—"
    else:
        razao_str = "—"
    _replace_on_slide(slide, "[X,Xx]x", razao_str)

    # Nota de soundness — usar justificativa completa do rating
    soundness = str(data["justificativa"])[:300] if data["justificativa"] != "—" else str(data["parecer"])
    _replace_on_slide(
        slide,
        "[Resumo da tese de crédito: por que os riscos estão mitigados. 2-3 linhas.]",
        soundness,
    )


# ---------------------------------------------------------------------------
# Fallback: gera apresentacao minima se template nao existir
# ---------------------------------------------------------------------------
def _generate_fallback(analise: dict, parametros: dict, output_path: str) -> str:
    """Gera apresentacao minimalista quando o template nao esta disponivel."""
    prs = Presentation()
    prs.slide_width = 9144000  # 10in
    prs.slide_height = 5143500  # 5.625in

    layout = prs.slide_layouts[6]  # blank

    # Slide unico com dados basicos
    slide = prs.slides.add_slide(layout)
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    tomador = parametros.get("tomador", _safe_get(analise, "tomador", "nome", default="N/A"))
    tipo = parametros.get("tipo_operacao", "N/A")
    volume = _fmt_brl(parametros.get("volume", 0))
    nota = _safe_get(analise, "rating_final", "nota", default="N/A")

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "ZYN CAPITAL — TEASER"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x22, 0x30, 0x40)

    for line in [
        f"\n\nTomador: {tomador}",
        f"Instrumento: {tipo}",
        f"Volume: {volume}",
        f"Rating: {nota}",
        f"\n{_current_date_pt()}",
        "\n\nCONFIDENCIAL — DISTRIBUIÇÃO RESTRITA",
        "\n⚠ Template oficial nao encontrado. Reinstale o template em templates/Teaser_Template_ZYN.pptx",
    ]:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = line
        run2.font.size = Pt(12)
        run2.font.color.rgb = RGBColor(0x34, 0x40, 0x50)

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# Funcao principal
# ---------------------------------------------------------------------------
def generate_teaser(
    analise: dict[str, Any],
    parametros: dict[str, Any],
    output_path: str,
) -> str:
    """
    Gera Teaser ZYN (.pptx) a partir do template oficial.

    Abre o template de 5 slides, substitui os placeholders pelos dados reais
    da operacao, e salva no output_path. Preserva toda a formatacao visual.

    Args:
        analise: Resultado da analise de credito (MAC).
        parametros: Parametros da operacao (tomador, volume, taxa, etc.).
        output_path: Caminho de saida do arquivo .pptx.

    Returns:
        Caminho do arquivo gerado.
    """
    # Garante que analise e parametros sao dicts
    if not isinstance(analise, dict):
        analise = {}
    if not isinstance(parametros, dict):
        parametros = {}

    if not TEMPLATE_PATH.exists():
        return _generate_fallback(analise, parametros, output_path)

    prs = Presentation(str(TEMPLATE_PATH))

    slides = list(prs.slides)
    if len(slides) < 5:
        return _generate_fallback(analise, parametros, output_path)

    _fill_cover(slides[0], analise, parametros)
    _fill_resumo(slides[1], analise, parametros)
    _fill_overview(slides[2], analise, parametros)
    _fill_estrutura(slides[3], analise, parametros)
    _fill_garantias(slides[4], analise, parametros)

    prs.save(output_path)
    return output_path
